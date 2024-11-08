from readline import insert_text
import streamlit as st 
import openai_helper3
import base64
from io import BytesIO
import PyPDF2
import docx
from pptx import Presentation
import requests 
from bs4 import BeautifulSoup
from openai import ChatCompletion

try:

# Utility function to load an image and convert it to base64
 def load_image(image_path):
    with open(image_path, "rb") as image_file:
        return base64.b64encode(image_file.read()).decode()

 def create_downloadable_business_case(extracted_data):
    business_case = f"Extracted Data:\n{extracted_data}"
    buffer = BytesIO()
    buffer.write(business_case.encode('utf-8'))
    buffer.seek(0)
    return buffer

 def clear_and_rerun():
    for key in st.session_state.keys():
        del st.session_state[key]
    st.experimental_rerun()

# Initialize session state variables
 if 'news_article' not in st.session_state:
    st.session_state.news_article = ""

 if 'extracted_data' not in st.session_state:
    st.session_state.extracted_data = ""

 if 'uploaded_file' not in st.session_state:
    st.session_state.uploaded_file = None

 if 'show_input' not in st.session_state:
    st.session_state.show_input = False

# Set theme color
 st.set_page_config(page_title="Acquia Value Selling Accelerator", page_icon="ac8.png", layout="wide", initial_sidebar_state="expanded")

# Path to the logo image
 logo_path = "/Users/taeeb.khan/Downloads/ac7.png"  # Update with your actual logo path
 logo_base64 = load_image(logo_path)

# Using HTML to create a colored title along with the logo
 html_title = f"""
    <div style="display: flex; align-items: start;">
        <img src="data:image/png;base64,{logo_base64}" alt="Logo" style="width: 65px; height: 65px; margin-right: 10px; margin-top: 0%">
        <h1 style='color:#1989C9; margin-top: -2.5%; font-size: 55px'>Acquia Value Selling Accelerator</h1>
    </div>
    """
 st.markdown(html_title, unsafe_allow_html=True)
 
 #st.markdown("<p style='margin-bottom: 0px; margin-top: 0px; font-size: 21px;'> >>> Extract three times for optimal results 🎯:</p>", unsafe_allow_html=True)
 st.markdown("<p style='margin-bottom: 0px; margin-top: 0px; font-size: 21px;'> >>> Please upload the file below to get started. Or else, press the button below to copy the content. (Extract three times for optimal results 🎯):</p>", unsafe_allow_html=True)
 
 def read_text_file(uploaded_file):
    return uploaded_file.read().decode('utf-8')

 def read_pdf_file(uploaded_file):
    reader = PyPDF2.PdfReader(uploaded_file)
    text = ""
    for page in reader.pages:
        text += page.extract_text()
    return text

 def read_docx_file(uploaded_file):
    doc = docx.Document(uploaded_file)
    return "\n".join([para.text for para in doc.paragraphs])

 def read_pptx_file(uploaded_file):
    prs = Presentation(uploaded_file)
    text = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text.append(shape.text)
    return "\n".join(text)
 
 
 # Create layout to align Drag & Drop and Copy-Paste side by side
 col1, col2 = st.columns([1, 1])

    # Sidebar with file uploader and copy-paste button
 with col1:

    uploaded_file = st.file_uploader("", type=['txt', 'pdf', 'docx', 'pptx'])
    st.markdown(
    """
    <style>
    .larger {
        margin-top: 0px;
        font-size: 20px; /* Adjust size as desired */
    }
    </style>
    """,
    unsafe_allow_html=True
    )    
    
 with col2:
    

    # Add a line above the copy-paste button
    st.write(" ")
    st.write(" ")
    st.markdown("<br/>" * -1, unsafe_allow_html=True)  # Creates 0 lines of vertical space
            
    if st.button("Copy-Paste"):
        st.session_state.show_input = True 
        
        
    st.markdown(
        """
        <style>
        .larger-text {
            margin-top: 20px;
            font-size: 20px; /* Adjust size as desired */
        }
        </style>
        """,
        unsafe_allow_html=True
        )   

 if uploaded_file is not None:
    file_type = uploaded_file.type
    if file_type == "text/plain":
        content = read_text_file(uploaded_file)
    elif file_type == "application/pdf":
        content = read_pdf_file(uploaded_file)
    elif file_type == "application/vnd.openxmlformats-officedocument.wordprocessingml.document":
        content = read_docx_file(uploaded_file)
    elif file_type == "application/vnd.openxmlformats-officedocument.presentationml.presentation":
        content = read_pptx_file(uploaded_file)
    else:
        st.error("Unsupported file type")
        content = ""
 
    st.success(f"File '{uploaded_file.name}' uploaded successfully!")
    st.session_state.news_article = content

# Conditionally display the news article text area
 if st.session_state.show_input or st.session_state.news_article:
    news_article = st.text_area("", value=st.session_state.get('news_article', ''), height=110, key='news_article')

 else:
    # Ensure that news_article is retrieved from session state if available
    news_article = st.session_state.get('news_article', '')
    
# Inject CSS to adjust spacing between success message and text area
 st.markdown(
     """
     <style>
        /* Target the success message container to reduce bottom margin */
        div.stAlert {
            margin-top: 10px;
            margin-bottom: -100px;
            width: fit-content;
        }
        div.stAlert p, div.stAlert div {
            margin: 0 !important;        /* Remove top/bottom margin */
            padding: 0 !important;       /* Remove padding */
        }
     </style>
     """,
    unsafe_allow_html=True
    )
 
# Display buttons in a row under the text area
 col1, col2, col3, col4, col5, col6, col7, col8 = st.columns([1, 1, 1, 1, 0.9, 0.75, 1.05, 1])
 with col5:
  if st.button("-- Discovery"):
        if news_article.strip():
            with st.spinner("Extracting Discovery Insights..."):
                 st.session_state.extracted_data = openai_helper3.extract_data(news_article)
        else:
                 st.warning("Report Missing!")
 with col6:
  if st.button("-- Email"):
        if news_article.strip():
            with st.spinner("Extracting Sample Email..."):
                 st.session_state.extracted_data = openai_helper3.extract_email_data(news_article)
        else:
                 st.warning("Report Missing!")
 with col7:
  if st.button("-- Business Case"):
        if news_article.strip():
            with st.spinner("Extracting Business Case..."):
                 st.session_state.extracted_data = openai_helper3.extract_business_data(news_article)
        else:
                 st.warning("Report Missing!")
 with col8:
  if st.button("-- Refresh",):
     clear_and_rerun()
 
 # Conditionally display the news article text area
 if st.session_state.show_input or st.session_state.news_article:
    st.markdown("<br/>" * 0, unsafe_allow_html=True)  # Creates 0 lines of vertical space
    st.markdown('<p style="font-size: 21px;">Summary:</p>', unsafe_allow_html=True)
    st.markdown("<div id='output' style='overflow-y: auto; font-family: sans-serif;font-size: 18px;'>{}</div>".format(st.session_state.extracted_data), unsafe_allow_html=True)

# Add a button for downloading the business case
 if st.session_state.extracted_data:
    business_case_file_buffer = create_downloadable_business_case(st.session_state.extracted_data)
    st.download_button(
        label="Download",
        data=business_case_file_buffer,
        file_name="business_case.txt",
        mime="text/plain"
    )


# Streamlit app layout
 st.markdown("<hr/>", unsafe_allow_html=True)
 # Conditionally display the news article text area
 if st.session_state.show_input or st.session_state.news_article:
    st.markdown("<div style='font-size: 21px;margin-bottom: -100px;'>Follow-up question concerning the submitted report:</div>", unsafe_allow_html=True)
    question = st.text_input("", "", key="question_text")

    # Button to submit the question
    if st.button("Get Answer"):
        if question.strip() and st.session_state.get('news_article', '').strip():
            with st.spinner("Looking for an answer..."):
                try:
                    st.session_state.answer = openai_helper3.extract_answer(st.session_state['news_article'], question)
                    st.markdown(f"<p><strong>Answer:</strong> {st.session_state.answer}</p>", unsafe_allow_html=True)
                except Exception as e:
                    st.error(f"Could not retrieve answer: {e}")
        else:
            st.warning("Please enter a question and ensure the content is provided.")

except Exception as e:
  st.error (f" An error occurred. Please try again later: {e}")

# Your main app goes here
st.markdown("<br/>" * 2, unsafe_allow_html=True)  # Creates 0 lines of vertical space
st.markdown('<p style="font-size: 21px;font-style: normal;">The Sales Enablement Team</p>', unsafe_allow_html=True)
